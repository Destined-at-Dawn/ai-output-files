#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
道法术器PPT生成器
功能：将Markdown格式的道法术器内容转换为PPT演示文稿
"""

import sys
import io
import re

# 修复Windows控制台编码问题
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding='utf-8')


def parse_markdown(md_content: str) -> dict:
    """解析Markdown内容"""
    result = {
        'title': '',
        'slides': []
    }

    lines = md_content.split('\n')
    current_slide = None
    current_content = []

    for line in lines:
        line = line.strip()

        # 解析标题
        if line.startswith('# ') and not result['title']:
            result['title'] = line[2:].strip()
            continue

        # 识别新幻灯片
        if line.startswith('## '):
            # 保存上一张幻灯片
            if current_slide:
                result['slides'].append({
                    'title': current_slide,
                    'content': current_content
                })

            current_slide = line[3:].strip()
            current_content = []
        elif current_slide:
            if line.startswith('-') or line.startswith('*'):
                content = line[1:].strip()
                if content.startswith('**'):
                    # 子标题
                    current_content.append({
                        'type': 'subheading',
                        'text': content.replace('**', '').strip()
                    })
                else:
                    current_content.append({
                        'type': 'bullet',
                        'text': content
                    })
            elif line and not line.startswith('>') and not line.startswith('---'):
                current_content.append({
                    'type': 'text',
                    'text': line
                })

    # 保存最后一张幻灯片
    if current_slide:
        result['slides'].append({
            'title': current_slide,
            'content': current_content
        })

    return result


def generate_ppt(parsed_data: dict, output_file: str):
    """生成PPT"""
    try:
        from pptx import Presentation
        from pptx.util import Pt, Inches
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor
    except ImportError:
        print("⚠️  需要安装python-pptx库")
        print("   请运行: pip install python-pptx")
        return False

    prs = Presentation()

    # 设置幻灯片尺寸为16:9
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # 标题页
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]

    title.text = parsed_data['title']
    subtitle.text = "道法术器拆解"

    # 内容页
    for slide_data in parsed_data['slides']:
        # 使用空白布局
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # 添加标题
        title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        title_frame = title.text_frame
        title_frame.text = slide_data['title']
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(139, 107, 76)  # 大地色

        # 添加内容
        y_pos = 1.3
        for item in slide_data['content']:
            if item['type'] == 'subheading':
                # 子标题
                textbox = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos), Inches(8.6), Inches(0.5))
                tf = textbox.text_frame
                tf.text = item['text']
                para = tf.paragraphs[0]
                para.font.size = Pt(18)
                para.font.bold = True
                para.font.color.rgb = RGBColor(212, 181, 158)
                y_pos += 0.6
            else:
                # 列表项
                textbox = slide.shapes.add_textbox(Inches(1), Inches(y_pos), Inches(8.3), Inches(0.4))
                tf = textbox.text_frame
                tf.text = "• " + item['text']
                para = tf.paragraphs[0]
                para.font.size = Pt(14)
                para.font.color.rgb = RGBColor(90, 74, 58)
                y_pos += 0.5

            # 每张幻灯片最多显示6-7项
            if y_pos > 4.5:
                break

    # 保存PPT
    prs.save(output_file)
    return True


def main():
    if len(sys.argv) < 2:
        print("用法: python generate_ppt.py <markdown文件路径>")
        print("示例: python generate_ppt.py transcript_cleaned.md")
        sys.exit(1)

    md_file = sys.argv[1]

    try:
        with open(md_file, 'r', encoding='utf-8') as f:
            md_content = f.read()

        # 解析Markdown
        parsed_data = parse_markdown(md_content)

        # 生成输出文件名
        output_file = md_file.replace('.md', '.pptx').replace('_cleaned', '')

        # 生成PPT
        if generate_ppt(parsed_data, output_file):
            print(f"✅ PPT生成完成！")
            print(f"📁 输出文件: {output_file}")
        else:
            print("❌ PPT生成失败")

    except Exception as e:
        print(f"❌ 错误: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)


if __name__ == '__main__':
    main()
